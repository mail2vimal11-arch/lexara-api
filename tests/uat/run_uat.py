"""
End-to-end UAT for the live LexAra deployment.

Drives the production frontend with Playwright (Chromium), captures a
screenshot per step, and emits a single .pptx report summarising pass /
fail per case.

Usage (locally or in CI):
    pip install -r tests/uat/requirements.txt
    playwright install --with-deps chromium
    python tests/uat/run_uat.py

Environment variables (all optional):
    UAT_BASE_URL          default https://lexara.tech
    UAT_API_BASE_URL      default https://api.lexara.tech
    UAT_OUTPUT_DIR        default qa-reports/<UTC-timestamp>
    UAT_TEST_USER_EMAIL   if set, runs the authenticated-flow case
    UAT_TEST_USER_PASSWORD
    UAT_HEADLESS          default 1; set 0 to watch the run locally

Exit code 0 on full pass, 1 if any case fails. The .pptx is written
regardless so failures are visible in the report.
"""

from __future__ import annotations

import os
import sys
import time
import traceback
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path

import httpx
from playwright.sync_api import Browser, Page, sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

BASE_URL = os.environ.get("UAT_BASE_URL", "https://lexara.tech").rstrip("/")
API_BASE_URL = os.environ.get("UAT_API_BASE_URL", "https://api.lexara.tech").rstrip("/")
HEADLESS = os.environ.get("UAT_HEADLESS", "1") != "0"
TEST_EMAIL = os.environ.get("UAT_TEST_USER_EMAIL")
TEST_PASSWORD = os.environ.get("UAT_TEST_USER_PASSWORD")

RUN_TIMESTAMP = datetime.now(timezone.utc).strftime("%Y%m%d-%H%M%SZ")
OUTPUT_DIR = Path(os.environ.get("UAT_OUTPUT_DIR", f"qa-reports/{RUN_TIMESTAMP}"))
SCREENSHOT_DIR = OUTPUT_DIR / "screenshots"


# ---------------------------------------------------------------------------
# Test result model
# ---------------------------------------------------------------------------

@dataclass
class StepResult:
    name: str
    status: str  # "pass" | "fail" | "skip"
    detail: str = ""
    screenshot: Path | None = None
    duration_ms: int = 0
    started_at: str = field(default_factory=lambda: datetime.now(timezone.utc).isoformat())


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _shot(page: Page, name: str) -> Path:
    """Save a full-page screenshot. Returns its path."""
    SCREENSHOT_DIR.mkdir(parents=True, exist_ok=True)
    safe = name.lower().replace(" ", "_").replace("/", "_")
    path = SCREENSHOT_DIR / f"{safe}.png"
    page.screenshot(path=str(path), full_page=True)
    return path


def _verify_page(page: Page, url: str, title_substring: str, name: str) -> StepResult:
    """Navigate, wait for load, assert title, screenshot."""
    started = time.monotonic()
    try:
        page.goto(url, wait_until="networkidle", timeout=30_000)
        actual = page.title()
        shot = _shot(page, name)
        if title_substring.lower() in actual.lower():
            return StepResult(
                name=name, status="pass",
                detail=f"GET {url} | <title> contains '{title_substring}' (got '{actual}')",
                screenshot=shot,
                duration_ms=int((time.monotonic() - started) * 1000),
            )
        return StepResult(
            name=name, status="fail",
            detail=f"GET {url} | expected title to contain '{title_substring}', got '{actual}'",
            screenshot=shot,
            duration_ms=int((time.monotonic() - started) * 1000),
        )
    except Exception as exc:
        # Best-effort screenshot even on failure
        try:
            shot = _shot(page, f"{name}_error")
        except Exception:
            shot = None
        return StepResult(
            name=name, status="fail",
            detail=f"Exception: {exc.__class__.__name__}: {exc}",
            screenshot=shot,
            duration_ms=int((time.monotonic() - started) * 1000),
        )


def _api_check(path: str, expect_keys: list[str], name: str) -> StepResult:
    """Hit an API endpoint and assert specific JSON keys are present."""
    started = time.monotonic()
    url = f"{API_BASE_URL}{path}"
    try:
        resp = httpx.get(url, timeout=15.0)
        body = resp.json() if "application/json" in resp.headers.get("content-type", "") else {}
        missing = [k for k in expect_keys if k not in body]
        if resp.status_code == 200 and not missing:
            preview = ", ".join(f"{k}={body[k]!r}" for k in expect_keys[:4])
            return StepResult(
                name=name, status="pass",
                detail=f"GET {url} -> 200 | {preview}",
                duration_ms=int((time.monotonic() - started) * 1000),
            )
        return StepResult(
            name=name, status="fail",
            detail=f"GET {url} -> {resp.status_code} | missing keys: {missing} | body={str(body)[:200]}",
            duration_ms=int((time.monotonic() - started) * 1000),
        )
    except Exception as exc:
        return StepResult(
            name=name, status="fail",
            detail=f"Exception: {exc.__class__.__name__}: {exc}",
            duration_ms=int((time.monotonic() - started) * 1000),
        )


# ---------------------------------------------------------------------------
# UAT cases
# ---------------------------------------------------------------------------

PAGE_CASES: list[tuple[str, str, str]] = [
    # (path, expected title substring, case name)
    ("/", "LexAra", "Landing page"),
    ("/auth.html", "Sign In", "Auth page"),
    ("/procurement.html", "Procurement Tools", "Procurement Tools page"),
    ("/procurement-ai.html", "Procurement AI", "Procurement AI page"),
    ("/procurement-intelligence.html", "Procurement Intelligence", "Procurement Intelligence page"),
    ("/negotiation-arena.html", "Negotiation Arena", "Negotiation Arena page"),
    ("/privacy.html", "Privacy Policy", "Privacy page"),
    ("/terms.html", "Terms of Service", "Terms page"),
]


def run_authenticated_flow(page: Page) -> list[StepResult]:
    """Optional: log in and verify the authenticated landing renders."""
    if not (TEST_EMAIL and TEST_PASSWORD):
        return [StepResult(
            name="Authenticated flow",
            status="skip",
            detail="UAT_TEST_USER_EMAIL / UAT_TEST_USER_PASSWORD not set — skipping",
        )]

    results: list[StepResult] = []
    started = time.monotonic()
    try:
        page.goto(f"{BASE_URL}/auth.html", wait_until="networkidle", timeout=30_000)
        # Selectors are pessimistic — we look for any input[type=email] / input[type=password].
        page.fill("input[type=email]", TEST_EMAIL)
        page.fill("input[type=password]", TEST_PASSWORD)
        page.click("button[type=submit], input[type=submit]")
        page.wait_for_load_state("networkidle", timeout=30_000)
        shot = _shot(page, "post_login")
        results.append(StepResult(
            name="Authenticated flow — sign in",
            status="pass" if "auth.html" not in page.url else "fail",
            detail=f"after submit: url={page.url}",
            screenshot=shot,
            duration_ms=int((time.monotonic() - started) * 1000),
        ))
    except Exception as exc:
        try:
            shot = _shot(page, "auth_flow_error")
        except Exception:
            shot = None
        results.append(StepResult(
            name="Authenticated flow — sign in",
            status="fail",
            detail=f"Exception: {exc.__class__.__name__}: {exc}",
            screenshot=shot,
            duration_ms=int((time.monotonic() - started) * 1000),
        ))
    return results


# ---------------------------------------------------------------------------
# Report (.pptx)
# ---------------------------------------------------------------------------

def write_pptx(results: list[StepResult], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank layout
    title_layout = prs.slide_layouts[0]  # title slide

    # ── Title slide ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = "LexAra — End-to-End UAT Report"
    subtitle = s.placeholders[1]
    passed = sum(1 for r in results if r.status == "pass")
    failed = sum(1 for r in results if r.status == "fail")
    skipped = sum(1 for r in results if r.status == "skip")
    subtitle.text = (
        f"Run: {RUN_TIMESTAMP}\n"
        f"Target: {BASE_URL}  |  API: {API_BASE_URL}\n"
        f"Pass: {passed}   Fail: {failed}   Skip: {skipped}"
    )

    # ── Summary slide ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    tbox = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = tbox.text_frame
    tf.text = "Summary"
    tf.paragraphs[0].runs[0].font.size = Pt(28)
    tf.paragraphs[0].runs[0].font.bold = True

    rows = len(results) + 1
    table = s.shapes.add_table(rows, 4, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8)).table
    headers = ["#", "Case", "Status", "Detail"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
    for idx, r in enumerate(results, start=1):
        table.cell(idx, 0).text = str(idx)
        table.cell(idx, 1).text = r.name
        table.cell(idx, 2).text = r.status.upper()
        table.cell(idx, 3).text = r.detail[:200]
        for col in range(4):
            for run in table.cell(idx, col).text_frame.paragraphs[0].runs:
                run.font.size = Pt(10)

    # ── One slide per result ──────────────────────────────────────────────
    for idx, r in enumerate(results, start=1):
        s = prs.slides.add_slide(blank)
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        tf = title.text_frame
        tf.text = f"{idx}. {r.name}  —  {r.status.upper()}"
        tf.paragraphs[0].runs[0].font.size = Pt(24)
        tf.paragraphs[0].runs[0].font.bold = True

        meta = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.5))
        meta.text_frame.text = f"{r.detail}   |   {r.duration_ms} ms   |   {r.started_at}"
        meta.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

        if r.screenshot and r.screenshot.exists():
            s.shapes.add_picture(
                str(r.screenshot),
                Inches(0.5), Inches(1.7),
                width=Inches(12.3),
            )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------

def main() -> int:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    results: list[StepResult] = []

    print(f"UAT run starting | base={BASE_URL} api={API_BASE_URL} headless={HEADLESS}")
    print(f"Output dir: {OUTPUT_DIR.resolve()}")

    # API health smoke (no browser needed)
    results.append(_api_check("/health", ["status"], "API /health"))
    results.append(_api_check("/status", ["status", "service", "version", "uptime_seconds"], "API /status"))

    # Browser flows
    with sync_playwright() as pw:
        browser: Browser = pw.chromium.launch(headless=HEADLESS)
        ctx = browser.new_context(viewport={"width": 1440, "height": 900})
        page = ctx.new_page()

        for path, expected_title, name in PAGE_CASES:
            print(f"  [{name}] -> {BASE_URL}{path}")
            results.append(_verify_page(page, f"{BASE_URL}{path}", expected_title, name))

        results.extend(run_authenticated_flow(page))

        ctx.close()
        browser.close()

    # Report
    pptx_path = OUTPUT_DIR / f"Lexara_UAT_{RUN_TIMESTAMP}.pptx"
    write_pptx(results, pptx_path)

    # Console summary
    print("\n=== UAT summary ===")
    for r in results:
        marker = {"pass": "PASS", "fail": "FAIL", "skip": "SKIP"}[r.status]
        print(f"  [{marker}] {r.name} — {r.detail[:120]}")
    print(f"\nReport: {pptx_path}")

    failed = sum(1 for r in results if r.status == "fail")
    return 1 if failed else 0


if __name__ == "__main__":
    try:
        sys.exit(main())
    except Exception:
        traceback.print_exc()
        sys.exit(2)
